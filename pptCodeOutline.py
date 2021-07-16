import glob
import os, sys
import string
import os.path
import sys
import shutil
#import xlrd
import csv
import argparse
import xlsxwriter
import pptx
import imghdr
#import antigravity
import pandas as pd
import numpy as np
import matplotlib as plt
import seaborn as sns
import tkinter as tk
from tkinter.filedialog import *
from tkinter.messagebox import *
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from datetime import date
from pathlib import Path
from PIL import Image
from io import StringIO, BytesIO
import cv2
from PIL import JpegImagePlugin


# Define random useful functions here


# Function Purpose: function changes a single character into a float.
# Precondition: s is a single character string.
# Postcondition: float(s) is the single character string changed to a float and s is the single character string returned as a string.


def maybe_float(s):
	try:
		return float(s)
	except (ValueError, TypeError):
		return s


# Function Purpose: function removes the extension at the end of a filename.
# Precondition: extensionlist is a full list of filenames containing the extension (.jpg, .png, .heic, etc...) at the end.
# Postcondition: filenames is a full list of filenames without the extension at the end.


def stem_list(extensionlist):
	noextensionlist = []

	for file in extensionlist:
		y = Path(file).stem
		noextensionlist.append(y)


	return noextensionlist
									

# Precondition: Filenames is a full list of image filenames as strings. Criteria is a user inputted string representing the beginning condition we want the filename to match.
# Postcondition: A list is returned containing only the filenames for modules matching the given criteria.
def sortcriteria(filenames, criteria):
	imagefiles = []
	errorfiles = []
	for path, dirs, files in os.walk('.'):
		for file in filenames:
			if len(file[:file.rfind("-")]) <= 16:
				firstpart = file[:file.rfind("-")]
				if criteria in firstpart:
					imagefiles.append(file)  
				else:
					errorfiles.append(file)
			elif len(file[:file.rfind("-")]) >= 16:
				firstpart = file[:(file[:file.rfind("-")]).rfind("-")]
				if criteria in firstpart:
					imagefiles.append(file)
				else:
					errorfiles.append(file)


	print("The following files don't fit the given criteria:", errorfiles)

	return imagefiles

def sortcriteriahelper(criterialist, pptx_location, prs):
	for i in criterialist:
		photos = []
		count = 0 
		k =0
		while count != 3:
			if i[k] == "-":
				count+=1
				k+=1
			else:
				k+=1
		newstring = i[:k]
		print(newstring)
		for p in criterialist:
			if newstring in p:
				photos.append(p)
		title = i[:k-1]
		addimages(photos, pptx_location, prs, title)





# Precondition: Filenames is a full list of image filenames as strings. Condition is a user inputted string representing the condition of the module(s) at the time the photo(s) were taken.
# Postcondition: A list is returned containing only the filenames for modules with the desired condition.


def sortcondition(filenames, condition):
	imagefiles = []
	errorfiles = []
	for path, dirs, files in os.walk('.'):
		for file in filenames:
			if len(file[:file.rfind("-")]) >= 16:
				postcondition = file[file.rfind("-")+1:file.rfind(".")]
				if str(postcondition) == str(condition):
					imagefiles.append(file)
				else:
					errorfiles.append(file)
			elif len(file[:file.rfind("-")]) <= 16:
				errorfiles.append(file)

	print("The following files don't fit the given postcondition:", errorfiles)

	return imagefiles


# Function purpose is to sort images into the powerpoint with each slide being of modules with different pre criterias.
# Precondition: criterialist is a full list of filenames with the same pre criteria as strings.
# Postcondition: the function adds a slide and inserts a slide and inserts all the images of the given pre criteria into the slide.


def formatbycriteria():
	#Creating this list just in case the number of images per slide will be more than 9 images.
		#POSSIBLE IMPROVEMENT AREA: Find some way to let the user know (maybe a pop up screen) when this happens.
	listtoobig = []

	#The user selected pre criteria from the GUI.
	criteria = criteriaentry.get()

	#Prompt the user to select the directory that contains all the possible image files they want to filter from.
	location = askdirectory(title = "Select location of image files")
	if(not location):
		sys.exit(0)

	#Change the current directory to the selected location.
	os.chdir(location)
	path = os.getcwd()
	#Get a list of all files in the current directory that are jpeg images.
	filenames = glob.glob("*.jpg")

	savename = savenameentry.get()

	pptx_location = askopenfilename(filetypes = [("Powerpoint Files", "*.pptx")], title = "Select Powerpoint Template File.")
	if not(pptx_location):
		sys.exit(0)

	prs = Presentation(pptx_location)

	#Find out how many photos have a module ID that matches the desired pre criteria.
	criterialist = sortcriteria(filenames, str(criteria))

	sortcriteriahelper(criterialist, pptx_location, prs)

	save_location = askdirectory(title = "Select location to save presentation.")
	os.chdir(save_location)
	prs.save(str(savename)+".pptx")


# Function purpose is to sort images into the powerpoint with each slide being of modules with different conditions.
# Precondition: conditionlist is a full list of filenames with the same condition as strings.
# Postcondition: the function adds a slide and inserts all the images of the given condition into the slide.


def formatbycondition():
	#Creating this list just in case the number of images per slide will be more than 9 images.
		#POSSIBLE IMPROVEMENT AREA: Find some way to let the user know (maybe a pop up screen) when this happens.
	listtoobig = []

	#The user selected pre criteria from the GUI.
	condition = conditionentry.get()

	#Prompt the user to select the directory that contains all the possible image files they want to filter from.
	location = askdirectory(title = "Select location of image files")
	if(not location):
		sys.exit(0)

	#Change the current directory to the selected location.
	os.chdir(location)
	#Get a list of all files in the current directory that are jpeg images.
	filenames = glob.glob("*.jpg")

	#Find out how many photos have a module ID that matches the desired post condition.
	conditionlist = sortcondition(filenames, str(condition))

	#Select the right number of images to insert per slide.
	addimages(conditionlist)


# Precondition: a full list of image filenames is passed into the function.
# Postcondition: the function adds the 1 image from the list into the slide.



def addimages(modulelist, pptx_location, prs,title):

	lyt_1 = prs.slide_layouts[0]
	lyt_2 = prs.slide_layouts[1]
	lyt_3 = prs.slide_layouts[2]
	lyt_4 = prs.slide_layouts[3]
	lyt_5 = prs.slide_layouts[4]
	lyt_6 = prs.slide_layouts[5]
	lyt_7 = prs.slide_layouts[6]
	lyt_8 = prs.slide_layouts[7]

	pictureslides = prs.slides.add_slide(lyt_6)
	pictureslidestitle = pictureslides.shapes.title


	pictureslidestitle.text = str(title)
	pictureslidestitle.text_frame.paragraphs[0].font.size = Pt(40)

	
	padding = 2
	for image in modulelist:
		JpegImagePlugin._getmp = lambda x: None
		im = Image.open(image)
		im = cv2.imread(image)

		image_width, image_height, color = im.shape
		slide_width = prs.slide_width.inches -4
		slide_height = prs.slide_height.inches -2
		if (image_width / slide_width) > (image_height / slide_height):
			#Image fits slide horizontally and must be scaled down vertically
		    left = Inches(5)
		    top = Inches(2.5)
		    pictureslides.shapes.add_picture(image, left, top, height = Inches(4))
		else:
		    # Image fits slide vertically and must be scaled down horizontally
			left = Inches(4)
			top = Inches(2.5)	
			pictureslides.shapes.add_picture(image, left, top, width = Inches(6))

		# pictureslides.shapes.add_picture(image, left_6_1, top_6_1, width_6, height_6)

	

# Precondition: function will take location of image files and then sort them into the separate lists.
# Postcondition: function will return a list of filenames all of the same criteria.


def getcriteriaimagefiles():
	location = askdirectory(title = "Select location of image files")
	if(not location):
		sys.exit(0)

	os.chdir(location)
	filenames = glob.glob("*.jpg")

	criterialist = sortcriteria(filenames, "EP002-P01")

	return criterialist


# Precondition: function will take location of image files and then sort them into the separate lists.
# Postcondition: function will return a list of filenames all of the same condition.


def getconditionimagefiles():
	location = askdirectory(title = "Select location of image files")
	if(not location):
		sys.exit(0) 

	os.chdir(location)
	filenames = glob.glob("*.jpg")

	conditionlist = sortcondition(filenames, "DH500")

	return conditionlist



# Part 3: Create the GUI


powerpointgui = tk.Tk()


# Create homepage screen.


powerpointcanvas = tk.Canvas(powerpointgui, width = 700, height = 600)
powerpointcanvas.pack()


# Create boxes to take user input.


criteriaentry = tk.Entry(powerpointgui)
powerpointcanvas.create_window(220, 100, window = criteriaentry)

criterialabel = tk.Label(text = "Pre Criteria")
powerpointcanvas.create_window(87.5, 100, window = criterialabel)


conditionentry = tk.Entry(powerpointgui)
powerpointcanvas.create_window(220, 140, window = conditionentry)

conditionlabel = tk.Label(text = "Post Condition")
powerpointcanvas.create_window(87.5, 140, window = conditionlabel)

savenameentry = tk.Entry(powerpointgui)
powerpointcanvas.create_window(520, 100, window = savenameentry)

savenamelabel = tk.Label(text = "Powerpoint File Name")
powerpointcanvas.create_window(380, 100, window = savenamelabel)


# Create buttons to select powerpoint formatting.


formatbyconditionbutton = tk.Button(text = "Format Powerpoint by Post Condition", command = formatbycondition)
powerpointcanvas.create_window(450, 140, window = formatbyconditionbutton)

formatbycriteriabutton = tk.Button(text = "Format Powerpoint by Pre Criteria", command = formatbycriteria)
powerpointcanvas.create_window(450, 180, window = formatbycriteriabutton)

powerpointgui.mainloop()
