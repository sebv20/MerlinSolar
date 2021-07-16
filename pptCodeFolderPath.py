import glob
import os, sys
import string
import os.path
import sys
import shutil
#import xlrd
import csv
import argparse
import xlsxwriter
import pptx
import imghdr
#import antigravity
import pandas as pd
import numpy as np
import matplotlib as plt
import seaborn as sns
import tkinter as tk
from tkinter.filedialog import *
from tkinter.messagebox import *
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from datetime import date
from pathlib import Path
from PIL import Image
from io import StringIO, BytesIO
import cv2
from PIL import JpegImagePlugin



# Part 1: How to import all the files we want


# Define random useful functions here


# Precondition: s is a single character string.
# Postcondition: float(s) is the single character string changed to a float and s is the single character string returned as a string.


def maybe_float(s):
	try:
		return float(s)
	except (ValueError, TypeError):
		return s


# Precondition: extensionlist is a full list of filenames containing the extension (.jpg, .png, .heic, etc...) at the end.
# Postcondition: filenames is a full list of filenames without the extension at the end.


def stem_list(extensionlist):
	noextensionlist = []

	for file in extensionlist:
		y = Path(file).stem
		noextensionlist.append(y)

	return noextensionlist

# Define useful script specific functions here

# Precondition: function is tied to a button input.
# Postcondition: function will take all the images from all the subdirectories within a directory and add them all onto a different slide per subdirectory.


def sorthelper(path, postcondition, pptx_location, prs):
	condition = os.scandir(path)

	for folder in condition:
		if folder.is_file():
			cwd = os.getcwd()
			imagefiles = []
			filenames = glob.glob("*.jpg")
			filenames_jpeg = glob.glob("*jpeg")

			for file in filenames:
				imagefiles.append(file)
			for file in filenames_jpeg:
				imagefiles.append(file)
			addimages(imagefiles, postcondition, pptx_location, prs, imagefiles,cwd)		
		

		else:
			os.chdir(folder)
			cwd = os.getcwd()
			if folders_in(cwd) == False:
				imagefiles = []
				filenames = glob.glob("*.jpg")
				filenames_jpeg = glob.glob("*jpeg")
				#filenames = glob.glob(cwd, recursive = True)
				for file in filenames:
					imagefiles.append(file)
				for file in filenames_jpeg:
					imagefiles.append(file)
				addimages(imagefiles, postcondition, pptx_location, prs, imagefiles,cwd)	
			else:
				sorthelper(cwd, postcondition, pptx_location, prs)
				# sorthelper(cwd, postcondition, pptx_location, prs)

def folders_in(path):
	for fname in os.listdir(path):
		if os.path.isdir(os.path.join(path, fname)):
			return True
	return False


def sortcondition():

	location = askdirectory(title = "Select post condition folder.")
	if(not location):
		sys.exit(0)

	splitlocationpath = os.path.split(location)
	postcondition = splitlocationpath[1]

	condition = os.scandir(location)

	savename = savenameentry.get()

	pptx_location = askopenfilename(filetypes = [("Powerpoint Files", "*.pptx")], title = "Select Powerpoint Template File.")
	if not(pptx_location):
		sys.exit(0)

	prs = Presentation(pptx_location) 

	# imagefiles = []
	# for root, subd, files in os.walk(location):
	# 	for name in files:
	# 		os.chdir(files)
	# 		#	file_name=glob.glob("*.jpg")
	# 		for file in file_name:
	# 			imagefiles.append(file)
	# 		addimages(subd, postcondition, pptx_location, prs, imagefiles)


	# for folder in condition:
	# 		if folder.is_dir():
	# 			imagefiles = []

	# 			os.chdir(folder)
	# 			filenames = glob.glob("*.jpg")

	# 			for file in filenames:
	# 				imagefiles.append(file)

	# 			addimages(folder, postcondition, pptx_location, prs, imagefiles)
	# 		else:
	# 			sorthelper(folder)
	
	sorthelper(location, postcondition, pptx_location, prs)
	save_location = askdirectory(title = "Select location to save presentation.")
	os.chdir(save_location)
	prs.save(str(savename)+".pptx")


# Precondition: pptx_location is the location of the powerpoint file template. prs is the powerpoint presentation. modulelist is the list of images from the given subdirectory.
# Postcondition: function will all all the images from a certain directory to a new slide.


def addimages(folder, postcondition, pptx_location, prs, modulelist, path):

	lyt_1 = prs.slide_layouts[0]
	lyt_2 = prs.slide_layouts[1]
	lyt_3 = prs.slide_layouts[2]
	lyt_4 = prs.slide_layouts[3]
	lyt_5 = prs.slide_layouts[4]
	lyt_6 = prs.slide_layouts[5]
	lyt_7 = prs.slide_layouts[6]
	lyt_8 = prs.slide_layouts[7]

	pictureslides = prs.slides.add_slide(lyt_6)
	pictureslidestitle = pictureslides.shapes.title
	pictureslidestitle.size = Pt(18)

	# newfolder = str(folder)
	# foldername = newfolder[newfolder.find("'")+1:newfolder.rfind("'")]


	pictureslidestitle.text = str(path)
	pictureslidestitle.text_frame.paragraphs[0].font.size = Pt(24)

	
	padding = 2
	for image in modulelist:
		JpegImagePlugin._getmp = lambda x: None
		im = Image.open(image)
		im = cv2.imread(image)

		image_width, image_height, color = im.shape
		slide_width = prs.slide_width.inches -4
		slide_height = prs.slide_height.inches -2
		if (image_width / slide_width) > (image_height / slide_height):
			#Image fits slide horizontally and must be scaled down vertically
		    print(1)
		    left = Inches(5)
		    top = Inches(2.5)
		    pictureslides.shapes.add_picture(image, left, top, height = Inches(4))
		else:
		    # Image fits slide vertically and must be scaled down horizontally
			print(2)
			left = Inches(4)
			top = Inches(2.5)	
			pictureslides.shapes.add_picture(image, left, top, width = Inches(6))
		# Convert from EMU to inches


# Part 3: Create the GUI


powerpointgui = tk.Tk()


# Create homepage screen.


powerpointcanvas = tk.Canvas(powerpointgui, width = 400, height = 220)
powerpointcanvas.pack()


# Create boxes to take user input.


savenameentry = tk.Entry(powerpointgui)
powerpointcanvas.create_window(200, 90, window = savenameentry)

savenamelabel = tk.Label(text = "Powerpoint File Name")
powerpointcanvas.create_window(200, 50, window = savenamelabel)


# Create buttons to select powerpoint formatting.


sortconditionbutton = tk.Button(text = "Format Powerpoint by Post Condition", command = sortcondition)
powerpointcanvas.create_window(200, 130, window = sortconditionbutton)

powerpointgui.mainloop()
